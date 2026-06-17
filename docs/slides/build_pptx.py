"""
Build the HOPE-MOVE overview deck as a native, fully-editable PowerPoint (.pptx).

Engine: python-pptx (+ Pillow for image sizing). Style is a restrained, branded
system (cream canvas, dark titles, a single blue accent, kicker labels, thin
accent rules, light-blue callouts, real screenshots + real result figures), built
to read as a designed deck, not a generic template.

Run:  python docs/slides/build_pptx.py
Out:  docs/slides/HOPE-MOVE-overview.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from PIL import Image

# ----------------------------------------------------------------------------- paths
HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
RESEARCH = os.path.join(ASSETS, "research")
OUT = os.path.join(HERE, "HOPE-MOVE-overview.pptx")
def A(n): return os.path.join(ASSETS, n)
def R(n): return os.path.join(RESEARCH, n)

# WEBP-safe loader (python-pptx rejects WEBP even with a .png name)
_NORM = os.path.join(ASSETS, ".build_norm"); os.makedirs(_NORM, exist_ok=True)
def safe_img(path):
    if not os.path.exists(path): return path
    try:
        with Image.open(path) as im:
            if (im.format or "").upper() in ("PNG","JPEG","BMP","GIF","TIFF"): return path
            out = os.path.join(_NORM, os.path.splitext(os.path.basename(path))[0]+".png")
            im.convert("RGBA").save(out, "PNG"); return out
    except Exception:
        return path

# ----------------------------------------------------------------------------- palette
CREAM      = RGBColor(0xFB,0xFA,0xF5)
WHITE      = RGBColor(0xFF,0xFF,0xFF)
INK        = RGBColor(0x17,0x1A,0x1F)
BODY       = RGBColor(0x3A,0x3D,0x42)
MUTED      = RGBColor(0x73,0x76,0x7C)
ACCENT     = RGBColor(0x1F,0x4E,0x9B)
ACCENT2    = RGBColor(0x2F,0x6F,0xD6)
CALLOUT_BG = RGBColor(0xE7,0xEF,0xFB)
CARD_BG    = WHITE
CARD_LINE  = RGBColor(0xE3,0xE2,0xD8)
HAIR       = RGBColor(0xDD,0xDD,0xD3)
BAR_GRAY   = RGBColor(0xC9,0xCE,0xD6)
FONT = "Segoe UI"
FONT_L = "Segoe UI Light"

# ----------------------------------------------------------------------------- deck
prs = Presentation()
prs.slide_width  = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
MARGIN = Inches(0.78); CW = SW - 2*MARGIN
_pg = {"n":0}

# ----------------------------------------------------------------------------- helpers
def _no_shadow(shape):
    sp = shape._element.spPr
    for tag in ("a:effectLst","a:effectDag"):
        for el in sp.findall(qn(tag)): sp.remove(el)
    sp.append(sp.makeelement(qn("a:effectLst"), {}))

def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    _pg["n"] += 1
    return s

def _runs(p, text, color, size, accent=ACCENT, base_bold=False):
    for i, seg in enumerate(text.split("**")):
        if not seg: continue
        r = p.add_run(); r.text = seg; r.font.name = FONT; r.font.size = Pt(size)
        if i % 2: r.font.bold = True; r.font.color.rgb = accent
        else:     r.font.bold = base_bold; r.font.color.rgb = color

def tbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf

def para(tf, text, size, color=BODY, bold=False, align=PP_ALIGN.LEFT,
         after=6, before=0, line=1.08, markup=False, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(after); p.space_before = Pt(before); p.line_spacing = line
    if markup: _runs(p, text, color, size, base_bold=bold)
    else:
        r = p.add_run(); r.text = text; r.font.name = FONT; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return p

def header(s, title, kicker=None):
    """Kicker (optional) + bold title + thin accent rule; the designed header."""
    y = Inches(0.62)
    if kicker:
        _, tf = tbox(s, MARGIN, Inches(0.5), CW, Inches(0.32))
        para(tf, kicker.upper(), 12.5, color=ACCENT, bold=True, first=True, after=0)
        y = Inches(0.92)
    _, tf = tbox(s, MARGIN, y, CW, Inches(0.7))
    para(tf, title, 30, color=INK, bold=True, first=True, after=0)
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y + Inches(0.62), Inches(0.62), Pt(3.2))
    rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT; rule.line.fill.background(); _no_shadow(rule)

def bullets(s, x, y, w, items, size=18, gap=11, color=BODY, line=1.13, h=Inches(4.6)):
    _, tf = tbox(s, x, y, w, h)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line; p.space_after = Pt(gap)
        rb = p.add_run(); rb.text = "•   "; rb.font.name = FONT; rb.font.size = Pt(size)
        rb.font.color.rgb = ACCENT; rb.font.bold = True
        _runs(p, it, color, size)

def callout(s, text, y, x=MARGIN, w=None, size=16.5, h=Inches(0.8)):
    w = w or CW
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = CALLOUT_BG; box.line.fill.background(); _no_shadow(box)
    try: box.adjustments[0] = 0.09
    except Exception: pass
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.26); tf.margin_top = tf.margin_bottom = Inches(0.06)
    para(tf, text, size, color=ACCENT, first=True, markup=True, after=0, line=1.06)
    return box

def card(s, x, y, w, h, eyebrow, heading, body, accent=ACCENT):
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    rect.fill.solid(); rect.fill.fore_color.rgb = CARD_BG
    rect.line.color.rgb = CARD_LINE; rect.line.width = Pt(1); _no_shadow(rect)
    try: rect.adjustments[0] = 0.045
    except Exception: pass
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y + Inches(0.16), Inches(0.075), h - Inches(0.32))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background(); _no_shadow(bar)
    tf = rect.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.22); tf.margin_bottom = Inches(0.18)
    para(tf, eyebrow.upper(), 12, color=MUTED, bold=True, first=True, after=6)
    para(tf, heading, 18.5, color=INK, bold=True, after=7, line=1.0)
    para(tf, body, 14.5, color=BODY, after=0, line=1.16, markup=True)
    return rect

def stat_card(s, x, y, w, h, big, label, sub=None):
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    rect.fill.solid(); rect.fill.fore_color.rgb = CALLOUT_BG; rect.line.fill.background(); _no_shadow(rect)
    try: rect.adjustments[0] = 0.06
    except Exception: pass
    tf = rect.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.32); tf.margin_right = Inches(0.28)
    para(tf, big, 46, color=ACCENT, bold=True, first=True, after=2, line=0.95)
    para(tf, label, 17, color=INK, bold=True, after=4 if sub else 0, line=1.05)
    if sub: para(tf, sub, 13, color=MUTED, after=0, line=1.12)
    return rect

def img_dims(path):
    with Image.open(path) as im: return im.size

def screenshot(s, path, x, y, w, h, align="center", valign="middle", border=True):
    """Fit image in box (aspect-preserved); optional thin border for a framed look."""
    if not os.path.exists(path):
        ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0xEE,0xEE,0xEE)
        ph.text_frame.text = "missing: " + os.path.basename(path); return
    iw, ih = img_dims(path); scale = min(w/iw, h/ih)
    nw, nh = int(iw*scale), int(ih*scale)
    nx = x + (w-nw)//2 if align=="center" else (x if align=="left" else x+(w-nw))
    ny = y + (h-nh)//2 if valign=="middle" else (y if valign=="top" else y+(h-nh))
    if border:
        bd = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, nx, ny, nw, nh)
        bd.fill.background(); bd.line.color.rgb = CARD_LINE; bd.line.width = Pt(1); _no_shadow(bd)
    s.shapes.add_picture(safe_img(path), nx, ny, nw, nh)

def caption(s, text, y, x=MARGIN, w=None, align=PP_ALIGN.CENTER, color=MUTED, size=12.5):
    w = w or CW
    _, tf = tbox(s, x, y, w, Inches(0.4))
    para(tf, text, size, color=color, first=True, align=align, after=0, markup=True)

def footer(s, dark=False):
    c = RGBColor(0xCF,0xD6,0xE6) if dark else RGBColor(0xB2,0xB4,0xB7)
    _, tf = tbox(s, MARGIN, SH-Inches(0.5), Inches(5), Inches(0.3))
    para(tf, "HOPE-MOVE", 9.5, color=c, first=True, after=0)
    _, tf = tbox(s, SW-Inches(0.95), SH-Inches(0.5), Inches(0.6), Inches(0.3))
    para(tf, str(_pg["n"]), 10, color=c, first=True, align=PP_ALIGN.RIGHT, after=0)

def logos_row(s, cy, height=Inches(0.6), gap=Inches(0.72)):
    specs = []
    for n in ["logo-innovateuk.png","logo-hope.png","logo-coventry.png"]:
        iw, ih = img_dims(A(n)); specs.append((A(n), int(height*iw/ih), int(height)))
    total = sum(w for _,w,_ in specs) + int(gap)*(len(specs)-1)
    x = int((SW-total)/2)
    for p,w,h in specs:
        s.shapes.add_picture(safe_img(p), x, int(cy-h/2), w, h); x += w+int(gap)

def add_bar_chart(s, x, y, cx, cy, title, cats, vals, hi, ymin=0.80, ymax=0.88, label_fs=9.0, title_fs=12.5):
    cd = CategoryChartData(); cd.categories = cats; cd.add_series("BERTScore", vals)
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, cd)
    ch = gf.chart; ch.has_legend = False; ch.has_title = True
    t = ch.chart_title.text_frame; t.text = title
    r0 = t.paragraphs[0].runs[0]
    r0.font.size = Pt(title_fs); r0.font.bold = True; r0.font.color.rgb = INK; r0.font.name = FONT
    plot = ch.plots[0]; plot.gap_width = 45; plot.has_data_labels = True
    dl = plot.data_labels; dl.number_format = "0.000"; dl.number_format_is_linked = False
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.size = Pt(label_fs); dl.font.name = FONT; dl.font.color.rgb = BODY
    va = ch.value_axis; va.minimum_scale = ymin; va.maximum_scale = ymax
    va.major_gridlines.format.line.color.rgb = RGBColor(0xEC,0xEC,0xE6)
    va.tick_labels.font.size = Pt(label_fs); va.tick_labels.font.name = FONT
    va.tick_labels.number_format = "0.00"; va.tick_labels.number_format_is_linked = False
    ca = ch.category_axis; ca.tick_labels.font.size = Pt(label_fs); ca.tick_labels.font.name = FONT
    ser = plot.series[0]; ser.format.fill.solid(); ser.format.fill.fore_color.rgb = BAR_GRAY
    for i in range(len(cats)):
        pt = ser.points[i]; pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = ACCENT2 if i == hi else BAR_GRAY
    return ch

# ============================================================================= SLIDES

# --- 1 · Title ---------------------------------------------------------------
s = slide(WHITE)
_, tf = tbox(s, Inches(1), Inches(1.95), SW-Inches(2), Inches(0.4))
para(tf, "COVENTRY UNIVERSITY  ·  HOPE-MOVE", 13.5, color=ACCENT, bold=True, first=True, align=PP_ALIGN.CENTER, after=0)
_, tf = tbox(s, Inches(1), Inches(2.45), SW-Inches(2), Inches(1.1))
para(tf, "AI Facilitator Assistant", 50, color=INK, bold=True, first=True, align=PP_ALIGN.CENTER, after=0)
rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(SW/2-Inches(0.5)), Inches(3.55), Inches(1.0), Pt(3.5))
rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT; rule.line.fill.background(); _no_shadow(rule)
_, tf = tbox(s, Inches(1.7), Inches(3.85), SW-Inches(3.4), Inches(1.0))
para(tf, "Helping facilitators spot who needs them, and reply with warmth, faster.",
     20, color=MUTED, first=True, align=PP_ALIGN.CENTER, after=0, line=1.2)
caption(s, "HOPE-MOVE  ·  People living with IIH 2025", Inches(4.75), align=PP_ALIGN.CENTER, size=13.5)
logos_row(s, Inches(6.05))
footer(s)

# --- 2 · Why we built it + what it does (merged) -----------------------------
s = slide()
header(s, "Why we built it and what it does")
_, tf = tbox(s, MARGIN, Inches(1.78), CW, Inches(1.1))
para(tf, "People living with long-term conditions can **quietly disengage** between sessions; about **1 in 4 drop out**, most going quiet in the **first week**. No facilitator can watch everyone, every week.",
     18, color=BODY, first=True, markup=True, after=0, line=1.28)
cy, ch = Inches(3.05), Inches(2.3); cw = (CW-Inches(0.5))/2
card(s, MARGIN, cy, cw, ch, "1 · Spot", "Who needs attention",
     "Every participant, **ranked each week**, with the plain-language reasons why.")
card(s, MARGIN+cw+Inches(0.5), cy, cw, ch, "2 · Support", "A warm draft reply",
     "A ready-to-send message in the programme's voice: **you review, edit, and send.**")
callout(s, "It **surfaces and suggests**. You **decide and send**.", Inches(5.6), h=Inches(0.8))
footer(s)

# --- 4 · The dashboard (hero) ------------------------------------------------
s = slide()
header(s, "One place to see your whole cohort")
screenshot(s, A("dashboard-hero.png"), MARGIN, Inches(1.78), CW, Inches(4.75), valign="top")
caption(s, "Follow-up queue  ·  participant detail  ·  AI outreach, in a single weekly view (live data).", Inches(6.75))
footer(s)

# --- 5 · SPOT (detail screenshot) --------------------------------------------
s = slide()
header(s, "See who needs you and why", kicker="Spot")
screenshot(s, A("dashboard-detail.png"), MARGIN, Inches(1.95), Inches(6.3), Inches(4.55), align="left", valign="top")
bx = MARGIN + Inches(6.7)
bullets(s, bx, Inches(2.05), SW-bx-MARGIN, [
    "Everyone **ranked each week**: Needs attention · Check in soon · On track.",
    "A clear **risk score** with the **reasons** behind it.",
    "A **recommended approach** and wellbeing cue for outreach.",
], size=18, gap=15, h=Inches(3.4))
callout(s, "Risk comes from a **Random Forest** model, shown in plain language.",
        Inches(5.5), x=bx, w=SW-bx-MARGIN, h=Inches(0.95), size=15)
footer(s)

# --- 6 · SUPPORT (outreach screenshot) ---------------------------------------
s = slide()
header(s, "Reply with warmth, faster", kicker="Support")
screenshot(s, A("dashboard-outreach.png"), MARGIN, Inches(1.9), Inches(3.05), Inches(4.7), align="left", valign="top")
bx = MARGIN + Inches(3.6)
bullets(s, bx, Inches(2.0), SW-bx-MARGIN, [
    "**Three warm tones**, drafted from the participant's own post.",
    "Written by a **fine-tuned small language model** (here, Qwen3).",
    "**Edit, polish, send**: always your words, your call.",
    "Your **feedback** (kept / edited / rejected) trains better drafts.",
], size=18, gap=15, h=Inches(3.6))
callout(s, "The draft is a **starting point**, never sent on its own.",
        Inches(5.85), x=bx, w=SW-bx-MARGIN, h=Inches(0.8), size=15.5)
footer(s)

# --- 6 · How it fits together (engines + architecture + safety, merged) ------
s = slide()
header(s, "How it fits together")
screenshot(s, A("architecture.png"), MARGIN, Inches(1.8), Inches(7.45), Inches(4.5), align="left", valign="middle", border=False)
bx = MARGIN + Inches(7.8)
bullets(s, bx, Inches(1.95), SW-bx-MARGIN, [
    "**Risk engine** (engagement_ml): predicts weekly dropout with a Random Forest, ~0.93–0.95 AUC.",
    "**Drafting engine** (comment_generation): a fine-tuned small language model trained on real Hope replies.",
    "**Safe by design:** anonymised · private (secrets server-side) · safety layer · nothing sent without you.",
], size=15.5, gap=15, h=Inches(3.7))
callout(s, "Both produce **suggestions only**; a **person always decides**.",
        Inches(6.0), x=bx, w=SW-bx-MARGIN, h=Inches(0.85), size=14)
footer(s)

# --- 10 · Section divider ----------------------------------------------------
s = slide(ACCENT)
_, tf = tbox(s, MARGIN, Inches(2.25), CW, Inches(0.4))
para(tf, "THE RESEARCH BEHIND IT", 14, color=RGBColor(0xB6,0xCD,0xF2), bold=True, first=True, after=0)
_, tf = tbox(s, MARGIN, Inches(2.7), CW, Inches(1.0))
para(tf, "Three research threads", 40, color=WHITE, bold=True, first=True, after=0)
_, tf = tbox(s, MARGIN, Inches(3.95), CW, Inches(1.0))
para(tf, "1.  Why retention happens   ·   2.  Predicting dropout   ·   3.  Generating replies",
     19, color=RGBColor(0xDD,0xE7,0xFA), first=True, after=0, line=1.3)
footer(s, dark=True)

# --- 11 · Evidence base ------------------------------------------------------
s = slide()
header(s, "What keeps people engaged", kicker="Evidence base · cohort study")
screenshot(s, R("te_dose_response.png"), MARGIN, Inches(2.15), Inches(6.7), Inches(4.3), align="left", valign="middle")
bx = MARGIN + Inches(7.05)
bullets(s, bx, Inches(2.2), SW-bx-MARGIN, [
    "**2,411 participants** across 7 health conditions.",
    "People who **write** complete at **89%** vs **48%** who don't (about **9× the odds**).",
    "More early activity → higher completion (a clear **dose-response**).",
    "But **727 browsed and never wrote**: exactly the gap we target.",
], size=17.5, gap=14, h=Inches(4.0))
caption(s, "Cohort study submitted to JMIR.", Inches(6.85), x=bx, w=SW-bx-MARGIN, align=PP_ALIGN.LEFT, size=12)
footer(s)

# --- 12 · Engagement profiles ------------------------------------------------
s = slide()
header(s, "Five engagement profiles", kicker="Evidence base · who to support first")
screenshot(s, R("te_profiles.png"), MARGIN, Inches(2.1), CW, Inches(3.6), valign="top")
bullets(s, MARGIN, Inches(5.85), CW, [
    "Drop-out ranges from **Disengaged (55%)** down to **Deep (0.3%)**.",
    "The **'Light' group** (browses but barely writes) is the **early-support target** the dashboard surfaces.",
], size=16.5, gap=7, h=Inches(1.2))
footer(s)

# --- Choosing the model (cold-start AUC across held-out courses/modules) -----
s = slide()
header(s, "Choosing the model", kicker="Risk engine · model selection")
screenshot(s, R("ml_cold_start_auc.png"), MARGIN, Inches(2.35), Inches(6.95), Inches(3.7), align="left", valign="middle")
bx = MARGIN + Inches(7.35)
bullets(s, bx, Inches(2.2), SW-bx-MARGIN, [
    "**Seven model families** compared (LR, RF, XGBoost, LightGBM, CatBoost, GRU, MLP).",
    "Tested **cold-start**: held out whole courses (LOCO) and modules (LOMO) the model never saw.",
    "The top models are **statistically tied** (~0.95 AUC); overlap alone can't separate them.",
    "**Random Forest** deployed for its robustness, calibration, and interpretability.",
], size=15.5, gap=13, h=Inches(4.0))
caption(s, "Cold-start AUC across held-out courses (LOCO) and modules (LOMO); the boxes overlap heavily.",
        Inches(6.5), x=MARGIN, w=Inches(7.0), align=PP_ALIGN.LEFT, size=12)
footer(s)

# --- Predicting dropout early (horizon AUC) ----------------------------------
s = slide()
header(s, "Predicting dropout early", kicker="Risk engine · performance")
screenshot(s, R("ml_horizon_auc.png"), MARGIN, Inches(2.15), Inches(6.7), Inches(4.3), align="left", valign="middle")
bx = MARGIN + Inches(7.05)
bullets(s, bx, Inches(2.2), SW-bx-MARGIN, [
    "**AUC 0.93 at week 1**, rising to ~**0.95**, catching dropout early.",
    "Stable across the **6-week** programme (a model per weekly horizon).",
    "Validated on **unseen courses and modules** (cold-start ≈ 0.95).",
    "Tuned for **high recall (≥90%)**, so few at-risk people are missed.",
], size=17, gap=14, h=Inches(4.0))
caption(s, "Submitted to J. of Biomedical Informatics.", Inches(6.85), x=bx, w=SW-bx-MARGIN, align=PP_ALIGN.LEFT, size=12)
footer(s)

# --- 11 · Drafting engine science (final forum roster chart) -----------------
s = slide()
header(s, "Drafting replies that sound like Hope", kicker="Drafting engine · small language models")
add_bar_chart(s, MARGIN, Inches(2.05), Inches(6.9), Inches(4.15),
    "Reply quality (BERTScore): forum-trained roster (n=1,228)",
    ["Qwen3.5-4B", "Qwen3-4B", "Qwen3-8B", "Qwen3-1.7B", "Qwen3-0.6B", "Gemma-4-E4B"],
    [0.8729, 0.8719, 0.8710, 0.8707, 0.8700, 0.8452], 0, label_fs=9.0, title_fs=11.5)
bx = MARGIN + Inches(7.25)
bullets(s, bx, Inches(2.0), SW-bx-MARGIN, [
    "**8 open models** fine-tuned with **QDoRA** on **3,900+ real Hope (post→reply) pairs**: programme data only, no synthetic augmentation.",
    "**Qwen3.5-4B** is the deployed default: top BERTScore (0.873).",
    "**Qwen3-4B** ties it within noise (~**0.002**) and has the **best GoalSetting fidelity**.",
    "**Safety layer:** crisis → MI policy → UK signposting · **human feedback** loop.",
], size=15, gap=11, h=Inches(3.7))
caption(s, "Also benchmarked on the activities sweep (n=771): SmolLM3-3B and Llama-3.2-3B as cross-family references. Both tables are in the NLP paper.",
        Inches(6.6), x=MARGIN, w=CW, align=PP_ALIGN.LEFT, size=11)
footer(s)

# --- 12 · Thank you ----------------------------------------------------------
s = slide(WHITE)
_, tf = tbox(s, Inches(1), Inches(2.3), SW-Inches(2), Inches(1.0))
para(tf, "Thank you", 46, color=INK, bold=True, first=True, align=PP_ALIGN.CENTER, after=0)
rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(SW/2-Inches(0.5)), Inches(3.4), Inches(1.0), Pt(3.5))
rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT; rule.line.fill.background(); _no_shadow(rule)
_, tf = tbox(s, Inches(1.6), Inches(3.7), SW-Inches(3.2), Inches(0.9))
para(tf, "HOPE-MOVE: AI to help facilitators keep people engaged and supported.",
     18.5, color=MUTED, first=True, align=PP_ALIGN.CENTER, after=0, line=1.2)
caption(s, "Michael Ajao-Olarinoye  ·  Coventry University  ·  funded by Innovate UK", Inches(4.7), align=PP_ALIGN.CENTER, size=13.5)
logos_row(s, Inches(6.05))
footer(s)

# ----------------------------------------------------------------------------- save
prs.save(OUT)
n_imgs = sum(1 for sl in prs.slides for sh in sl.shapes if sh.shape_type == 13)
print(f"Saved {OUT}")
print(f"Slides: {len(prs.slides._sldIdLst)}  |  embedded images: {n_imgs}")
